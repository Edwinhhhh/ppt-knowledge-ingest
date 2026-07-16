#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import fnmatch
import html
import hashlib
import json
import re
import shutil
import subprocess
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "dcterms": "http://purl.org/dc/terms/",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class SlideImage:
    path: str
    content_type: str
    width: int | None
    height: int | None


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def safe_name(text: str) -> str:
    text = re.sub(r"[\\/:*?\"<>|]+", "-", text).strip()
    return re.sub(r"\s+", " ", text)[:120] or "deck"


def ext_from_content_type(content_type: str) -> str:
    mapping = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/jpg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/svg+xml": ".svg",
        "image/x-emf": ".emf",
        "image/x-wmf": ".wmf",
    }
    return mapping.get(content_type, ".bin")


def extract_core_properties(pptx_path: Path) -> dict[str, str]:
    props: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        if "docProps/core.xml" not in zf.namelist():
            return props
        root = ET.fromstring(zf.read("docProps/core.xml"))
    for key in ["title", "subject", "creator", "description", "keywords", "lastModifiedBy"]:
        node = root.find(f"dc:{key}", NS)
        if node is None:
            node = root.find(f"cp:{key}", NS)
        if node is not None and node.text:
            props[key] = node.text.strip()
    for key in ["created", "modified"]:
        node = root.find(f"dcterms:{key}", NS)
        if node is not None and node.text:
            props[key] = node.text.strip()
    return props


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def emu(value: Any) -> int:
    try:
        return int(value)
    except Exception:
        return 0


def bbox_for_shape(shape: Any, slide_width: int, slide_height: int) -> dict[str, float | int]:
    left = emu(getattr(shape, "left", 0))
    top = emu(getattr(shape, "top", 0))
    width = emu(getattr(shape, "width", 0))
    height = emu(getattr(shape, "height", 0))
    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "right": left + width,
        "bottom": top + height,
        "left_pct": round(left / slide_width, 4) if slide_width else 0,
        "top_pct": round(top / slide_height, 4) if slide_height else 0,
        "width_pct": round(width / slide_width, 4) if slide_width else 0,
        "height_pct": round(height / slide_height, 4) if slide_height else 0,
    }


def font_rgb(run: Any) -> str | None:
    try:
        rgb = run.font.color.rgb
        return f"#{rgb}" if rgb is not None else None
    except Exception:
        return None


def run_hyperlink(run: Any) -> str | None:
    try:
        return run.hyperlink.address
    except Exception:
        return None


def extract_sections_by_slide(pptx_path: Path) -> tuple[dict[int, str], list[dict[str, Any]]]:
    by_slide: dict[int, str] = {}
    sections: list[dict[str, Any]] = []
    with zipfile.ZipFile(pptx_path) as zf:
        if "ppt/presentation.xml" not in zf.namelist():
            return by_slide, sections
        root = ET.fromstring(zf.read("ppt/presentation.xml"))

    slide_id_to_index: dict[str, int] = {}
    sld_id_lst = root.find("p:sldIdLst", NS)
    if sld_id_lst is not None:
        for idx, node in enumerate(list(sld_id_lst), start=1):
            slide_id = node.attrib.get("id")
            if slide_id:
                slide_id_to_index[slide_id] = idx

    for section in [node for node in root.iter() if local_name(node.tag) == "section"]:
        name = section.attrib.get("name") or "Untitled Section"
        slide_nums: list[int] = []
        for node in section.iter():
            if local_name(node.tag) != "sldId":
                continue
            slide_id = node.attrib.get("id")
            if slide_id and slide_id in slide_id_to_index:
                slide_no = slide_id_to_index[slide_id]
                by_slide[slide_no] = name
                slide_nums.append(slide_no)
        sections.append({"name": name, "slides": slide_nums})
    return by_slide, sections


def extract_notes_by_slide(pptx_path: Path) -> dict[int, list[str]]:
    notes: dict[int, list[str]] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        note_files = sorted(
            [name for name in zf.namelist() if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", name)],
            key=lambda x: int(re.search(r"notesSlide(\d+)\.xml", x).group(1)),
        )
        for name in note_files:
            slide_num = int(re.search(r"notesSlide(\d+)\.xml", name).group(1))
            root = ET.fromstring(zf.read(name))
            texts = []
            for node in root.findall(".//a:t", NS):
                if node.text and node.text.strip():
                    texts.append(node.text.strip())
            if texts:
                notes[slide_num] = texts
    return notes


def extract_text(shape: Any) -> list[str]:
    return [paragraph["text"] for paragraph in extract_paragraphs(shape)]


def extract_paragraphs(shape: Any) -> list[dict[str, Any]]:
    paragraphs: list[dict[str, Any]] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            runs = []
            for run in paragraph.runs:
                if not run.text:
                    continue
                runs.append(
                    {
                        "text": run.text,
                        "bold": bool(run.font.bold) if run.font.bold is not None else None,
                        "italic": bool(run.font.italic) if run.font.italic is not None else None,
                        "underline": bool(run.font.underline) if run.font.underline is not None else None,
                        "color": font_rgb(run),
                        "hyperlink": run_hyperlink(run),
                    }
                )
            text = "".join(run["text"] for run in runs).strip()
            if text:
                paragraphs.append({"text": text, "level": paragraph.level, "runs": runs})
    return paragraphs


def extract_text_block(shape: Any, slide_width: int, slide_height: int) -> dict[str, Any] | None:
    paragraphs = extract_paragraphs(shape)
    if not paragraphs:
        return None
    return {
        "kind": "text",
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "bbox": bbox_for_shape(shape, slide_width, slide_height),
        "text": "\n".join(paragraph["text"] for paragraph in paragraphs),
        "paragraphs": paragraphs,
    }


def extract_table(shape: Any) -> list[list[str]]:
    rows: list[list[str]] = []
    if not getattr(shape, "has_table", False):
        return rows
    for row in shape.table.rows:
        values = []
        for cell in row.cells:
            if getattr(cell, "is_spanned", False):
                values.append("")
            else:
                values.append(cell.text.strip())
        rows.append(values)
    return rows


def iter_categories(chart: Any) -> list[str]:
    try:
        categories = chart.plots[0].categories
    except Exception:
        return []
    values: list[str] = []
    for category in categories:
        label = getattr(category, "label", None)
        values.append(str(label if label is not None else category))
    return values


def chart_values(series: Any) -> list[Any]:
    try:
        return [value for value in series.values]
    except Exception:
        return []


def extract_chart_summary(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_chart", False):
        return None
    chart = shape.chart
    summary: dict[str, Any] = {"chart_type": str(chart.chart_type), "categories": iter_categories(chart)}
    errors = []
    try:
        if chart.has_title and chart.chart_title.text_frame.text:
            summary["title"] = chart.chart_title.text_frame.text.strip()
    except Exception as exc:
        errors.append(f"title: {exc}")
    try:
        summary["series"] = [
            {
                "name": series.name,
                "values": chart_values(series),
            }
            for series in chart.series
        ]
    except Exception as exc:
        errors.append(f"series: {exc}")
    if errors:
        summary["extraction_errors"] = errors
    return summary


def detect_special_object(shape: Any, slide_width: int, slide_height: int) -> dict[str, Any] | None:
    shape_type = getattr(shape, "shape_type", None)
    shape_type_name = getattr(shape_type, "name", str(shape_type))
    xml = getattr(getattr(shape, "element", None), "xml", "")
    detected: str | None = None
    if "drawingml/2006/diagram" in xml or "<dgm:" in xml:
        detected = "smartart"
    elif "OLE_OBJECT" in shape_type_name or "oleObj" in xml:
        detected = "embedded_ole_object"
    if not detected:
        return None
    return {
        "kind": detected,
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "shape_type": shape_type_name,
        "bbox": bbox_for_shape(shape, slide_width, slide_height),
        "note": "Detected but not deeply parsed; use slide rendering/OCR or OOXML-specific parsing for full fidelity.",
    }


def extract_slide_images(slide: Any, slide_no: int, images_dir: Path) -> list[SlideImage]:
    images: list[SlideImage] = []
    img_idx = 0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        img_idx += 1
        try:
            image = shape.image
        except ValueError:
            continue
        ext = ext_from_content_type(image.content_type)
        out_name = f"slide-{slide_no:03d}-image-{img_idx:02d}{ext}"
        out_path = images_dir / out_name
        out_path.write_bytes(image.blob)
        images.append(
            SlideImage(
                path=f"assets/images/{out_name}",
                content_type=image.content_type,
                width=getattr(image, "size", (None, None))[0],
                height=getattr(image, "size", (None, None))[1],
            )
        )
    return images


def load_tag_keywords(path: Path | None) -> list[str]:
    if path is None:
        return []
    if not path.exists():
        raise FileNotFoundError(path)
    if path.suffix.lower() == ".json":
        data = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(data, list):
            raise ValueError("Tag keyword JSON must be a list of strings.")
        return [str(item).strip() for item in data if str(item).strip()]
    return [
        line.strip()
        for line in path.read_text(encoding="utf-8").splitlines()
        if line.strip() and not line.strip().startswith("#")
    ]


def infer_tags(text: str, tag_keywords: list[str]) -> list[str]:
    return [tag for tag in tag_keywords if tag.lower() in text.lower()]


def first_nonempty(values: list[str]) -> str:
    for value in values:
        if value.strip():
            return value.strip()
    return ""


def layout_groups(text_blocks: list[dict[str, Any]], slide_height: int) -> list[dict[str, Any]]:
    if not text_blocks:
        return []
    groups: list[list[dict[str, Any]]] = []
    tolerance = max(int(slide_height * 0.035), 250000)
    for block in sorted(text_blocks, key=lambda b: (b["bbox"]["top"], b["bbox"]["left"])):
        center_y = block["bbox"]["top"] + block["bbox"]["height"] / 2
        placed = False
        for group in groups:
            group_center = sum(item["bbox"]["top"] + item["bbox"]["height"] / 2 for item in group) / len(group)
            if abs(center_y - group_center) <= tolerance:
                group.append(block)
                placed = True
                break
        if not placed:
            groups.append([block])

    result = []
    for idx, group in enumerate(groups, start=1):
        ordered = sorted(group, key=lambda b: b["bbox"]["left"])
        result.append(
            {
                "group_no": idx,
                "layout": "columns" if len(ordered) > 1 else "single",
                "block_ids": [item["shape_id"] for item in ordered],
                "texts": [item["text"] for item in ordered],
            }
        )
    return result


def build_slide_records(pptx_path: Path, output_dir: Path, tag_keywords: list[str] | None = None) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    tag_keywords = tag_keywords or []
    prs = Presentation(str(pptx_path))
    notes = extract_notes_by_slide(pptx_path)
    sections_by_slide, sections = extract_sections_by_slide(pptx_path)
    images_dir = output_dir / "assets" / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    slide_width = emu(prs.slide_width)
    slide_height = emu(prs.slide_height)

    records: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_blocks: list[str] = []
        rich_text_blocks: list[dict[str, Any]] = []
        tables: list[list[list[str]]] = []
        charts: list[dict[str, Any]] = []
        special_objects: list[dict[str, Any]] = []

        shapes = sorted(slide.shapes, key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)))
        for shape in shapes:
            text_block = extract_text_block(shape, slide_width, slide_height)
            if text_block:
                rich_text_blocks.append(text_block)
                text_blocks.extend(paragraph["text"] for paragraph in text_block["paragraphs"])
            table = extract_table(shape)
            if table:
                tables.append(table)
            chart = extract_chart_summary(shape)
            if chart:
                charts.append(chart)
            special_object = detect_special_object(shape, slide_width, slide_height)
            if special_object:
                special_objects.append(special_object)

        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        if not title:
            title = first_nonempty(text_blocks)
        if re.fullmatch(r"\s*\d+\s*", title):
            title = first_nonempty([block for block in text_blocks if not re.fullmatch(r"\s*\d+\s*", block)])

        body_text = "\n".join(text_blocks)
        note_text = "\n".join(notes.get(idx, []))
        all_text = "\n".join([title, body_text, note_text])
        records.append(
            {
                "slide_no": idx,
                "section": sections_by_slide.get(idx),
                "title": title or f"Slide {idx}",
                "text_blocks": text_blocks,
                "rich_text_blocks": rich_text_blocks,
                "layout_groups": layout_groups(rich_text_blocks, slide_height),
                "tables": tables,
                "charts": charts,
                "notes": notes.get(idx, []),
                "images": [image.__dict__ for image in extract_slide_images(slide, idx, images_dir)],
                "special_objects": special_objects,
                "tags": infer_tags(all_text, tag_keywords),
            }
        )
    suppress_repeated_notes(records)
    return records, sections


def suppress_repeated_notes(records: list[dict[str, Any]]) -> None:
    freq: dict[str, int] = {}
    for record in records:
        for note in record["notes"]:
            freq[note] = freq.get(note, 0) + 1
    threshold = max(3, len(records) // 2)
    for record in records:
        record["notes"] = [
            note
            for note in record["notes"]
            if freq.get(note, 0) < threshold and not re.fullmatch(r"[‹<#>\d./\-\s]+", note)
        ]


def markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    lines = [
        "| " + " | ".join(cell.replace("\n", " ") for cell in header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in normalized[1:]:
        lines.append("| " + " | ".join(cell.replace("\n", " ") for cell in row) + " |")
    return "\n".join(lines)


def markdown_rich_text(blocks: list[dict[str, Any]]) -> list[str]:
    lines: list[str] = []
    for block in blocks:
        for paragraph in block["paragraphs"]:
            indent = "  " * int(paragraph.get("level") or 0)
            lines.append(f"{indent}- {paragraph['text']}")
    return lines


def one_line(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def layout_text(record: dict[str, Any]) -> str:
    chunks: list[str] = []
    for group in record.get("layout_groups", []):
        if group["layout"] == "columns":
            chunks.append(
                " | ".join(f"Column {idx}: {one_line(text)}" for idx, text in enumerate(group["texts"], start=1))
            )
        else:
            chunks.extend(group["texts"])
    return "\n".join(chunks)


def write_markdown(path: Path, deck_title: str, metadata: dict[str, Any], records: list[dict[str, Any]]) -> None:
    lines = [
        f"# {deck_title}",
        "",
        "## Document Metadata",
        "",
        f"- Source file: `{metadata['source_file']}`",
        f"- Slide count: {metadata['slide_count']}",
        f"- SHA-256: `{metadata['sha256']}`",
        f"- Generated at: {metadata['generated_at']}",
        "",
        "## Slides",
        "",
    ]
    for record in records:
        lines.extend(
            [
                f"### Slide {record['slide_no']}: {record['title']}",
                "",
                "**核心文本**",
                "",
            ]
        )
        if record.get("section"):
            lines.extend([f"章节：`{record['section']}`", ""])

        if record["rich_text_blocks"]:
            lines.extend(markdown_rich_text(record["rich_text_blocks"]))
        else:
            lines.append("- （未抽取到可编辑文本，可能为图片页或复杂图形页）")
        lines.append("")

        column_groups = [group for group in record.get("layout_groups", []) if group["layout"] == "columns"]
        if column_groups:
            lines.append("**版面结构**")
            lines.append("")
            for group in column_groups:
                lines.append(f"- 水平并列组 {group['group_no']}:")
                for col_idx, text in enumerate(group["texts"], start=1):
                    lines.append(f"  - Column {col_idx}: {one_line(text)}")
            lines.append("")

        if record["tables"]:
            lines.append("**表格**")
            lines.append("")
            for table_idx, table in enumerate(record["tables"], start=1):
                lines.append(f"表格 {table_idx}:")
                lines.append("")
                lines.append(markdown_table(table))
                lines.append("")

        if record["charts"]:
            lines.append("**图表对象**")
            lines.append("")
            for chart in record["charts"]:
                lines.append(f"- {json.dumps(chart, ensure_ascii=False)}")
            lines.append("")

        if record.get("special_objects"):
            lines.append("**特殊对象**")
            lines.append("")
            for item in record["special_objects"]:
                lines.append(f"- {item['kind']}: {item.get('name') or item.get('shape_id')}（{item['note']}）")
            lines.append("")

        if record["notes"]:
            lines.append("**备注**")
            lines.append("")
            for item in record["notes"]:
                lines.append(f"- {item}")
            lines.append("")

        if record["images"]:
            lines.append("**图片资产**")
            lines.append("")
            for image in record["images"]:
                lines.append(f"- `{image['path']}` ({image['content_type']})")
            lines.append("")

        lines.extend(
            [
                "**标签**",
                "",
                ", ".join(f"`{tag}`" for tag in record["tags"]) if record["tags"] else "（待人工或模型补充）",
                "",
                "**来源**",
                "",
                f"`source_slide: {record['slide_no']}`",
                "",
            ]
        )
    path.write_text("\n".join(lines), encoding="utf-8")


def write_chunks(path: Path, deck_id: str, records: list[dict[str, Any]], source_file: str) -> None:
    with path.open("w", encoding="utf-8") as fh:
        for record in records:
            chunk_text = "\n".join(record["text_blocks"])
            structured_layout = layout_text(record)
            if structured_layout and structured_layout != chunk_text:
                chunk_text += "\n\n版面结构:\n" + structured_layout
            if record["charts"]:
                chunk_text += "\n\n图表数据:\n" + json.dumps(record["charts"], ensure_ascii=False)
            if record["notes"]:
                chunk_text += "\n\n备注:\n" + "\n".join(record["notes"])
            chunk = {
                "id": f"{deck_id}#slide-{record['slide_no']:03d}",
                "deck_id": deck_id,
                "slide_no": record["slide_no"],
                "source_file": source_file,
                "source_slide": record["slide_no"],
                "section": record.get("section"),
                "title": record["title"],
                "text": chunk_text.strip(),
                "layout_groups": record.get("layout_groups", []),
                "charts": record["charts"],
                "special_objects": record.get("special_objects", []),
                "tags": record["tags"],
                "source": {"type": "pptx", "file": source_file, "slide": record["slide_no"]},
                "assets": {"images": [image["path"] for image in record["images"]]},
            }
            fh.write(json.dumps(chunk, ensure_ascii=False) + "\n")


def read_jsonl(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    rows = []
    with path.open("r", encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if line:
                rows.append(json.loads(line))
    return rows


def write_jsonl(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as fh:
        for row in rows:
            fh.write(json.dumps(row, ensure_ascii=False) + "\n")


def emit_progress(enabled: bool, event: str, **payload: Any) -> None:
    if not enabled:
        return
    try:
        print(json.dumps({"event": event, **payload}, ensure_ascii=False), flush=True)
    except BrokenPipeError:
        return


def append_jsonl(path: Path, row: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("a", encoding="utf-8") as fh:
        fh.write(json.dumps(row, ensure_ascii=False) + "\n")


def html_list(items: list[str]) -> str:
    if not items:
        return "<p class=\"muted\">未抽取到可编辑文本，可能为图片页或复杂图形页。</p>"
    return "<ul>" + "".join(f"<li>{html.escape(item)}</li>" for item in items) + "</ul>"


def html_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    head = "".join(f"<th>{html.escape(cell)}</th>" for cell in header)
    body = "".join(
        "<tr>" + "".join(f"<td>{html.escape(cell)}</td>" for cell in row) + "</tr>"
        for row in normalized[1:]
    )
    return f"<table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table>"


def html_rich_text(blocks: list[dict[str, Any]]) -> str:
    if not blocks:
        return "<p class=\"muted\">未抽取到可编辑文本，可能为图片页或复杂图形页。</p>"
    items = []
    for block in blocks:
        for paragraph in block["paragraphs"]:
            level = int(paragraph.get("level") or 0)
            items.append(f"<li style=\"margin-left:{level * 18}px\">{html.escape(paragraph['text'])}</li>")
    return "<ul>" + "".join(items) + "</ul>"


def html_layout_groups(groups: list[dict[str, Any]]) -> str:
    column_groups = [group for group in groups if group["layout"] == "columns"]
    if not column_groups:
        return ""
    rendered = []
    for group in column_groups:
        columns = "".join(f"<div>{html.escape(text)}</div>" for text in group["texts"])
        rendered.append(f"<div class=\"layout-row\"><strong>水平并列组 {group['group_no']}</strong><div class=\"layout-cols\">{columns}</div></div>")
    return "<section class=\"block\"><h3>版面结构</h3>" + "".join(rendered) + "</section>"


def image_src_for_html(image: dict[str, Any], html_path: Path, standalone: bool) -> str:
    rel_path = image["path"]
    if not standalone:
        return rel_path
    image_path = html_path.parent / rel_path
    try:
        data = base64.b64encode(image_path.read_bytes()).decode("ascii")
        return f"data:{image['content_type']};base64,{data}"
    except Exception:
        return rel_path


def write_html(
    path: Path,
    deck_title: str,
    metadata: dict[str, Any],
    records: list[dict[str, Any]],
    standalone: bool = False,
) -> None:
    nav = "\n".join(
        f"<a href=\"#slide-{record['slide_no']:03d}\"><span>{record['slide_no']:02d}</span>{html.escape(record['title'])}</a>"
        for record in records
    )
    cards: list[str] = []
    for record in records:
        tags = "".join(f"<span class=\"tag\">{html.escape(tag)}</span>" for tag in record["tags"])
        if not tags:
            tags = "<span class=\"tag muted-tag\">待补充标签</span>"

        images = ""
        if record["images"]:
            images = "<div class=\"image-grid\">" + "".join(
                f"<figure><img src=\"{html.escape(image_src_for_html(image, path, standalone))}\" alt=\"Slide {record['slide_no']} image\"><figcaption>{html.escape(image['content_type'])}</figcaption></figure>"
                for image in record["images"]
            ) + "</div>"

        tables = ""
        if record["tables"]:
            tables = "<section class=\"block\"><h3>表格</h3>" + "".join(html_table(table) for table in record["tables"]) + "</section>"

        charts = ""
        if record["charts"]:
            chart_items = "".join(f"<li><code>{html.escape(json.dumps(chart, ensure_ascii=False))}</code></li>" for chart in record["charts"])
            charts = f"<section class=\"block\"><h3>图表对象</h3><ul>{chart_items}</ul></section>"

        notes = ""
        if record["notes"]:
            notes = f"<section class=\"block\"><h3>备注</h3>{html_list(record['notes'])}</section>"

        layout = html_layout_groups(record.get("layout_groups", []))
        special_objects = ""
        if record.get("special_objects"):
            items = "".join(
                f"<li>{html.escape(item['kind'])}: {html.escape(str(item.get('name') or item.get('shape_id')))} <span class=\"muted\">{html.escape(item['note'])}</span></li>"
                for item in record["special_objects"]
            )
            special_objects = f"<section class=\"block\"><h3>特殊对象</h3><ul>{items}</ul></section>"

        cards.append(
            f"""
            <article class="slide-card" id="slide-{record['slide_no']:03d}">
              <div class="slide-kicker">Slide {record['slide_no']:02d}</div>
              <h2>{html.escape(record['title'])}</h2>
              {f'<p class="section-label">章节：{html.escape(record["section"])}</p>' if record.get("section") else ''}
              <section class="block">
                <h3>核心文本</h3>
                {html_rich_text(record['rich_text_blocks'])}
              </section>
              {layout}
              {tables}
              {charts}
              {special_objects}
              {notes}
              {images}
              <div class="card-footer">
                <div class="tags">{tags}</div>
                <code>source_slide: {record['slide_no']}</code>
              </div>
            </article>
            """
        )

    page = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(deck_title)}</title>
  <style>
    :root {{
      --ink: #172026;
      --muted: #66737f;
      --line: #dfe5e8;
      --paper: #f6f3ed;
      --panel: #ffffff;
      --accent: #0f7c80;
      --accent-2: #b24d3e;
      --soft: #e8f4f3;
      --code: #28333a;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "PingFang SC", "Microsoft YaHei", sans-serif;
      color: var(--ink);
      background: var(--paper);
      line-height: 1.6;
    }}
    .app {{
      display: grid;
      grid-template-columns: 300px minmax(0, 1fr);
      min-height: 100vh;
    }}
    aside {{
      position: sticky;
      top: 0;
      height: 100vh;
      overflow: auto;
      padding: 28px 20px;
      border-right: 1px solid var(--line);
      background: #fbfaf7;
    }}
    main {{
      padding: 40px clamp(24px, 4vw, 64px);
    }}
    h1 {{
      margin: 0 0 10px;
      font-size: 32px;
      line-height: 1.2;
      letter-spacing: 0;
    }}
    .summary {{
      display: flex;
      flex-wrap: wrap;
      gap: 10px;
      margin: 18px 0 28px;
      color: var(--muted);
    }}
    .summary span {{
      border: 1px solid var(--line);
      background: var(--panel);
      padding: 6px 10px;
      border-radius: 6px;
    }}
    nav a {{
      display: grid;
      grid-template-columns: 36px 1fr;
      gap: 8px;
      align-items: start;
      padding: 9px 8px;
      color: var(--ink);
      text-decoration: none;
      border-radius: 6px;
      font-size: 14px;
    }}
    nav a:hover {{ background: var(--soft); }}
    nav span {{
      color: var(--accent);
      font-weight: 700;
    }}
    .slide-card {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 26px;
      margin: 0 0 22px;
      box-shadow: 0 10px 24px rgba(23, 32, 38, 0.05);
    }}
    .slide-kicker {{
      color: var(--accent-2);
      font-size: 13px;
      font-weight: 700;
      text-transform: uppercase;
      margin-bottom: 6px;
    }}
    h2 {{
      margin: 0 0 20px;
      font-size: 24px;
      line-height: 1.3;
      letter-spacing: 0;
    }}
    h3 {{
      margin: 0 0 8px;
      font-size: 15px;
      color: var(--muted);
    }}
    ul {{ margin: 0; padding-left: 20px; }}
    li + li {{ margin-top: 5px; }}
    .block {{ margin-top: 18px; }}
    .muted {{ color: var(--muted); }}
    .section-label {{
      margin: -8px 0 16px;
      color: var(--muted);
      font-size: 13px;
    }}
    .layout-row {{
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 12px;
      background: #fbfcfb;
    }}
    .layout-row + .layout-row {{ margin-top: 10px; }}
    .layout-cols {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(160px, 1fr));
      gap: 10px;
      margin-top: 8px;
    }}
    .layout-cols div {{
      border-left: 3px solid var(--accent);
      padding: 8px 10px;
      background: #fff;
    }}
    .image-grid {{
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
      gap: 14px;
      margin-top: 20px;
    }}
    figure {{
      margin: 0;
      border: 1px solid var(--line);
      border-radius: 8px;
      overflow: hidden;
      background: #f8faf9;
    }}
    img {{
      display: block;
      width: 100%;
      max-height: 260px;
      object-fit: contain;
      background: #fff;
    }}
    figcaption {{
      padding: 7px 9px;
      color: var(--muted);
      font-size: 12px;
      border-top: 1px solid var(--line);
    }}
    table {{
      width: 100%;
      border-collapse: collapse;
      margin-top: 8px;
      font-size: 14px;
    }}
    th, td {{
      border: 1px solid var(--line);
      padding: 8px;
      vertical-align: top;
    }}
    th {{ background: var(--soft); text-align: left; }}
    .card-footer {{
      display: flex;
      justify-content: space-between;
      gap: 16px;
      align-items: center;
      margin-top: 22px;
      padding-top: 16px;
      border-top: 1px solid var(--line);
    }}
    .tags {{
      display: flex;
      flex-wrap: wrap;
      gap: 6px;
    }}
    .tag {{
      display: inline-flex;
      align-items: center;
      min-height: 24px;
      padding: 2px 8px;
      border-radius: 999px;
      background: var(--soft);
      color: var(--accent);
      font-size: 12px;
      font-weight: 650;
    }}
    .muted-tag {{ color: var(--muted); background: #f0f1ef; }}
    code {{
      color: var(--code);
      background: #f3f5f4;
      border: 1px solid var(--line);
      border-radius: 5px;
      padding: 2px 5px;
      font-size: 12px;
      word-break: break-all;
    }}
    @media (max-width: 860px) {{
      .app {{ display: block; }}
      aside {{
        position: relative;
        height: auto;
        border-right: 0;
        border-bottom: 1px solid var(--line);
      }}
      main {{ padding: 24px 16px; }}
      .slide-card {{ padding: 20px; }}
      .card-footer {{ display: block; }}
      .card-footer code {{ display: inline-block; margin-top: 12px; }}
    }}
  </style>
</head>
<body>
  <div class="app">
    <aside>
      <h1>{html.escape(deck_title)}</h1>
      <div class="summary">
        <span>{metadata['slide_count']} 页</span>
        <span>HTML 试点版</span>
      </div>
      <nav>{nav}</nav>
    </aside>
    <main>
      {''.join(cards)}
    </main>
  </div>
</body>
</html>
"""
    path.write_text(page, encoding="utf-8")


def find_required_binary(name: str) -> str:
    path = shutil.which(name)
    if not path:
        raise RuntimeError(f"Required binary not found: {name}")
    return path


def run_checked(args: list[str]) -> None:
    try:
        subprocess.run(
            args,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except subprocess.CalledProcessError as exc:
        message = "\n".join(
            part
            for part in [
                f"Command failed: {' '.join(args)}",
                f"Exit code: {exc.returncode}",
                f"stdout: {exc.stdout.strip()}" if exc.stdout else "",
                f"stderr: {exc.stderr.strip()}" if exc.stderr else "",
            ]
            if part
        )
        raise RuntimeError(message) from exc


def render_pdf_to_slide_images(pdf_path: Path, output_dir: Path, dpi: int = 160) -> list[Path]:
    pdftoppm = find_required_binary("pdftoppm")
    render_dir = output_dir / "assets" / "rendered_slides"
    render_dir.mkdir(parents=True, exist_ok=True)
    for old in render_dir.glob("slide-*.png"):
        old.unlink()
    prefix = render_dir / "slide"
    run_checked([pdftoppm, "-r", str(dpi), "-png", str(pdf_path), str(prefix)])
    images = sorted(render_dir.glob("slide-*.png"))
    if not images:
        raise RuntimeError("pdftoppm did not produce slide images.")
    return images


def render_pptx_to_slide_images(pptx_path: Path, output_dir: Path, dpi: int = 160) -> tuple[Path, list[Path]]:
    soffice = find_required_binary("soffice")
    tmp_dir = output_dir / "tmp_render"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    run_checked(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(pptx_path)],
    )
    pdf_path = tmp_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        pdf_candidates = sorted(tmp_dir.glob("*.pdf"))
        if not pdf_candidates:
            raise RuntimeError("soffice did not produce a PDF.")
        pdf_path = pdf_candidates[0]

    images = render_pdf_to_slide_images(pdf_path, output_dir, dpi)
    return pdf_path, images


def rendered_img_src(image_path: Path, html_path: Path, standalone: bool) -> str:
    if standalone:
        data = base64.b64encode(image_path.read_bytes()).decode("ascii")
        return f"data:image/png;base64,{data}"
    return image_path.relative_to(html_path.parent).as_posix()


def write_rendered_html(
    path: Path,
    deck_title: str,
    slide_images: list[Path],
    slide_width: int,
    slide_height: int,
    standalone: bool = False,
) -> None:
    aspect = f"{slide_width} / {slide_height}" if slide_width and slide_height else "16 / 9"
    nav = "\n".join(
        f"<a href=\"#slide-{idx:03d}\">{idx:02d}</a>" for idx in range(1, len(slide_images) + 1)
    )
    slides = "\n".join(
        f"""
        <section class="slide" id="slide-{idx:03d}">
          <img src="{html.escape(rendered_img_src(image_path, path, standalone))}" alt="Slide {idx}">
        </section>
        """
        for idx, image_path in enumerate(slide_images, start=1)
    )
    page = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(deck_title)}</title>
  <style>
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: #202124;
      color: #f5f5f5;
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "PingFang SC", "Microsoft YaHei", sans-serif;
    }}
    header {{
      position: sticky;
      top: 0;
      z-index: 10;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 10px 16px;
      background: rgba(20, 22, 24, 0.92);
      border-bottom: 1px solid rgba(255,255,255,0.12);
      backdrop-filter: blur(12px);
    }}
    h1 {{
      margin: 0;
      font-size: 15px;
      font-weight: 650;
      letter-spacing: 0;
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
    }}
    nav {{
      display: flex;
      gap: 6px;
      overflow-x: auto;
      max-width: 50vw;
      padding-bottom: 2px;
    }}
    nav a {{
      color: #e8eaed;
      text-decoration: none;
      font-size: 12px;
      border: 1px solid rgba(255,255,255,0.18);
      border-radius: 6px;
      padding: 3px 7px;
    }}
    nav a:hover {{ background: rgba(255,255,255,0.12); }}
    main {{
      width: min(100%, 1280px);
      margin: 0 auto;
      padding: 22px clamp(10px, 2vw, 28px) 48px;
    }}
    .slide {{
      width: 100%;
      aspect-ratio: {aspect};
      margin: 0 auto 24px;
      background: #fff;
      box-shadow: 0 16px 48px rgba(0,0,0,0.38);
    }}
    .slide img {{
      display: block;
      width: 100%;
      height: 100%;
      object-fit: contain;
    }}
    @media print {{
      header {{ display: none; }}
      body {{ background: #fff; }}
      main {{ width: 100%; padding: 0; }}
      .slide {{ margin: 0; box-shadow: none; break-after: page; }}
    }}
  </style>
</head>
<body>
  <header>
    <h1>{html.escape(deck_title)}</h1>
    <nav>{nav}</nav>
  </header>
  <main>{slides}</main>
</body>
</html>
"""
    path.write_text(page, encoding="utf-8")


def display_path(path: Path, base: Path) -> str:
    try:
        return str(path.relative_to(base))
    except ValueError:
        return str(path)


def process_deck(
    pptx_path: Path,
    output_dir: Path,
    archive_source: bool = False,
    tag_keywords: list[str] | None = None,
    source_sha256: str | None = None,
    standalone_html: bool = False,
    rendered_html: bool = False,
    rendered_standalone_html: bool = False,
    rendered_source_pdf: Path | None = None,
) -> dict[str, Any]:
    pptx_path = pptx_path.resolve()
    output_dir = output_dir.resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, got: {pptx_path}")

    output_dir.mkdir(parents=True, exist_ok=True)

    core_props = extract_core_properties(pptx_path)
    records, sections = build_slide_records(pptx_path, output_dir, tag_keywords)
    deck_title = (core_props.get("title") or records[0]["title"]) if records else safe_name(pptx_path.stem)
    deck_id = safe_name(deck_title)
    prs = Presentation(str(pptx_path))
    slide_width = emu(prs.slide_width)
    slide_height = emu(prs.slide_height)

    source_sha256 = source_sha256 or sha256_file(pptx_path)
    metadata = {
        "deck_id": deck_id,
        "title": deck_title,
        "source_file": str(pptx_path),
        "original_copy": None,
        "file_size_bytes": pptx_path.stat().st_size,
        "sha256": source_sha256,
        "slide_count": len(records),
        "sections": sections,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "core_properties": core_props,
        "tagging": {
            "mode": "keyword_file" if tag_keywords else "none",
            "keyword_count": len(tag_keywords or []),
        },
        "outputs": {
            "markdown": "content.md",
            "html": "content.html",
            "records": "slides.json",
            "chunks": "chunks.jsonl",
            "embedded_images": "assets/images/",
        },
        "warnings": [],
        "limitations": [
            "This extracts editable PPT text, paragraph levels, approximate layout groups, notes, tables, chart data when accessible, and embedded images.",
            "Image-only slide text requires OCR or vision extraction in a later pass.",
            "SmartArt and OLE objects are detected but not fully parsed.",
        ],
    }

    if archive_source:
        archive_dir = output_dir / "source"
        archive_dir.mkdir(parents=True, exist_ok=True)
        shutil.copy2(pptx_path, archive_dir / "original.pptx")
        metadata["original_copy"] = "source/original.pptx"

    write_markdown(output_dir / "content.md", deck_title, metadata, records)
    write_html(output_dir / "content.html", deck_title, metadata, records)
    if standalone_html:
        metadata["outputs"]["standalone_html"] = "content.standalone.html"
        write_html(output_dir / "content.standalone.html", deck_title, metadata, records, standalone=True)
    (output_dir / "metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    (output_dir / "slides.json").write_text(json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8")
    write_chunks(output_dir / "chunks.jsonl", deck_id, records, str(pptx_path))
    if rendered_html or rendered_standalone_html:
        try:
            if rendered_source_pdf:
                pdf_path = rendered_source_pdf.resolve()
                if not pdf_path.exists():
                    raise FileNotFoundError(pdf_path)
                slide_images = render_pdf_to_slide_images(pdf_path, output_dir)
            else:
                pdf_path, slide_images = render_pptx_to_slide_images(pptx_path, output_dir)
            metadata["outputs"]["rendered_pdf"] = display_path(pdf_path, output_dir)
            if rendered_html:
                metadata["outputs"]["rendered_html"] = "rendered.html"
                write_rendered_html(output_dir / "rendered.html", deck_title, slide_images, slide_width, slide_height)
            if rendered_standalone_html:
                metadata["outputs"]["rendered_standalone_html"] = "rendered.standalone.html"
                write_rendered_html(output_dir / "rendered.standalone.html", deck_title, slide_images, slide_width, slide_height, standalone=True)
        except Exception as exc:
            metadata["warnings"].append(
                {
                    "type": "rendered_html_skipped",
                    "message": "Visual HTML was requested but skipped because slide rendering failed. Markdown, HTML, JSON, and chunks were still generated.",
                    "error": str(exc),
                }
            )
            metadata["limitations"].append(
                "Rendered slide screenshots were not generated in this run. Install LibreOffice/soffice and Poppler/pdftoppm, or provide a rendered-source PDF."
            )
        (output_dir / "metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")

    manifest = {
        "metadata": "metadata.json",
        "content": "content.md",
        "html": "content.html",
        "slides": "slides.json",
        "chunks": "chunks.jsonl",
    }
    if standalone_html:
        manifest["standalone_html"] = "content.standalone.html"
    if (output_dir / "rendered.html").exists():
        manifest["rendered_html"] = "rendered.html"
    if (output_dir / "rendered.standalone.html").exists():
        manifest["rendered_standalone_html"] = "rendered.standalone.html"
    if archive_source:
        manifest["source"] = "source/original.pptx"
    (output_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    return {"output_dir": str(output_dir), "slide_count": len(records), "title": deck_title, "sha256": source_sha256}


def is_excluded(path: Path, root: Path, patterns: list[str]) -> bool:
    if not patterns:
        return False
    rel = path.relative_to(root).as_posix()
    parts = path.relative_to(root).parts
    for pattern in patterns:
        if fnmatch.fnmatch(rel, pattern) or fnmatch.fnmatch(path.name, pattern):
            return True
        if any(fnmatch.fnmatch(part, pattern) for part in parts):
            return True
    return False


def iter_pptx_files(root: Path, exclude_patterns: list[str] | None = None) -> list[Path]:
    exclude_patterns = exclude_patterns or []
    return sorted(
        path
        for path in root.rglob("*")
        if path.is_file()
        and path.suffix.lower() == ".pptx"
        and not path.name.startswith("~$")
        and not is_excluded(path, root, exclude_patterns)
    )


def library_deck_dir(output_root: Path, source_sha256: str, pptx_path: Path) -> Path:
    slug = safe_name(pptx_path.stem).replace(" ", "-")
    return output_root / "decks" / f"{source_sha256[:16]}-{slug[:72]}"


def catalog_row(pptx_path: Path, source_sha256: str, status: str, **extra: Any) -> dict[str, Any]:
    stat = pptx_path.stat()
    row = {
        "path": str(pptx_path),
        "filename": pptx_path.name,
        "size_bytes": stat.st_size,
        "modified_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
        "sha256": source_sha256,
        "status": status,
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    row.update(extra)
    return row


def rebuild_all_chunks(output_root: Path, deck_dirs: list[Path]) -> int:
    all_chunks_path = output_root / "all_chunks.jsonl"
    count = 0
    seen_dirs: set[Path] = set()
    all_chunks_path.parent.mkdir(parents=True, exist_ok=True)
    with all_chunks_path.open("w", encoding="utf-8") as out:
        for deck_dir in deck_dirs:
            deck_dir = deck_dir.resolve()
            if deck_dir in seen_dirs:
                continue
            seen_dirs.add(deck_dir)
            chunks_path = deck_dir / "chunks.jsonl"
            metadata_path = deck_dir / "metadata.json"
            if not chunks_path.exists():
                continue
            metadata = {}
            if metadata_path.exists():
                metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
            for chunk in read_jsonl(chunks_path):
                chunk["deck_dir"] = str(deck_dir)
                if metadata.get("source_file"):
                    chunk["source_file"] = metadata["source_file"]
                    if isinstance(chunk.get("source"), dict):
                        chunk["source"].setdefault("file", metadata["source_file"])
                if chunk.get("slide_no") is not None:
                    chunk["source_slide"] = chunk["slide_no"]
                    if isinstance(chunk.get("source"), dict):
                        chunk["source"].setdefault("slide", chunk["slide_no"])
                if metadata.get("sha256"):
                    chunk["source_sha256"] = metadata["sha256"]
                out.write(json.dumps(chunk, ensure_ascii=False) + "\n")
                count += 1
    return count


def process_library(
    library_root: Path,
    output_root: Path,
    archive_source: bool = False,
    tag_keywords: list[str] | None = None,
    resume: bool = True,
    skip_duplicates: bool = True,
    dry_run: bool = False,
    exclude_patterns: list[str] | None = None,
    standalone_html: bool = False,
    rendered_html: bool = False,
    rendered_standalone_html: bool = False,
    progress_jsonl: bool = False,
) -> dict[str, Any]:
    library_root = library_root.resolve()
    output_root = output_root.resolve()
    if not library_root.is_dir():
        raise NotADirectoryError(library_root)
    output_root.mkdir(parents=True, exist_ok=True)

    catalog_path = output_root / "catalog.jsonl"
    failed_path = output_root / "failed.jsonl"
    decks: dict[str, Path] = {}
    rows: list[dict[str, Any]] = []
    failures: list[dict[str, Any]] = []
    processed_or_existing_dirs: list[Path] = []

    pptx_files = iter_pptx_files(library_root, exclude_patterns)
    total = len(pptx_files)
    emit_progress(progress_jsonl, "progress_start", total=total, current=0, status="running", phase="scan" if dry_run else "convert")

    for index, pptx_path in enumerate(pptx_files, 1):
        emit_progress(progress_jsonl, "file_start", total=total, current=index, status="running", file=str(pptx_path))
        try:
            source_sha256 = sha256_file(pptx_path)
            output_dir = library_deck_dir(output_root, source_sha256, pptx_path)
            if skip_duplicates and source_sha256 in decks:
                row = catalog_row(
                    pptx_path,
                    source_sha256,
                    "duplicate",
                    duplicate_of=str(decks[source_sha256]),
                    output_dir=str(decks[source_sha256]),
                )
                rows.append(row)
                emit_progress(progress_jsonl, "file_done", total=total, current=index, status="duplicate", file=str(pptx_path), output_dir=str(decks[source_sha256]))
                continue

            decks[source_sha256] = output_dir
            if not dry_run:
                processed_or_existing_dirs.append(output_dir)
            if dry_run:
                status = "skipped_existing" if resume and (output_dir / "manifest.json").exists() else "pending"
                rows.append(catalog_row(pptx_path, source_sha256, status, output_dir=str(output_dir)))
                emit_progress(progress_jsonl, "file_done", total=total, current=index, status=status, file=str(pptx_path), output_dir=str(output_dir))
                continue
            if resume and (output_dir / "manifest.json").exists():
                rows.append(catalog_row(pptx_path, source_sha256, "skipped_existing", output_dir=str(output_dir)))
                emit_progress(progress_jsonl, "file_done", total=total, current=index, status="skipped_existing", file=str(pptx_path), output_dir=str(output_dir))
                continue

            result = process_deck(
                pptx_path,
                output_dir,
                archive_source,
                tag_keywords,
                source_sha256,
                standalone_html,
                rendered_html,
                rendered_standalone_html,
            )
            rows.append(catalog_row(pptx_path, source_sha256, "processed", output_dir=str(output_dir), title=result["title"], slide_count=result["slide_count"]))
            emit_progress(progress_jsonl, "file_done", total=total, current=index, status="processed", file=str(pptx_path), output_dir=str(output_dir))
        except Exception as exc:
            source_sha256 = ""
            try:
                source_sha256 = sha256_file(pptx_path)
            except Exception:
                pass
            failure = catalog_row(pptx_path, source_sha256, "failed", error=str(exc))
            rows.append(failure)
            failures.append(failure)
            emit_progress(progress_jsonl, "file_done", total=total, current=index, status="failed", file=str(pptx_path), error=str(exc))

    write_jsonl(catalog_path, rows)
    write_jsonl(failed_path, failures)
    if dry_run:
        write_jsonl(output_root / "all_chunks.jsonl", [])
        chunk_count = 0
    else:
        chunk_count = rebuild_all_chunks(output_root, processed_or_existing_dirs)

    summary = {
        "library_root": str(library_root),
        "output_root": str(output_root),
        "dry_run": dry_run,
        "exclude_patterns": exclude_patterns or [],
        "pptx_count": len(rows),
        "pending_count": sum(1 for row in rows if row["status"] == "pending"),
        "processed_count": sum(1 for row in rows if row["status"] == "processed"),
        "skipped_existing_count": sum(1 for row in rows if row["status"] == "skipped_existing"),
        "duplicate_count": sum(1 for row in rows if row["status"] == "duplicate"),
        "failed_count": len(failures),
        "all_chunks_count": chunk_count,
        "catalog": str(catalog_path),
        "failed": str(failed_path),
        "all_chunks": str(output_root / "all_chunks.jsonl"),
    }
    (output_root / "library_manifest.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    emit_progress(progress_jsonl, "progress_complete", total=total, current=total, status="done", phase="scan" if dry_run else "convert")
    return summary


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert PPTX decks into AI-ready knowledge assets.")
    inputs = parser.add_mutually_exclusive_group(required=True)
    inputs.add_argument("--pptx", type=Path, help="Single PPTX file to convert.")
    inputs.add_argument("--input-dir", type=Path, help="Folder containing PPTX files to convert.")
    inputs.add_argument("--library-root", type=Path, help="Recursively scan a local PPTX library and build a resumable knowledge library.")
    parser.add_argument("--output-dir", type=Path, help="Output folder for a single PPTX conversion.")
    parser.add_argument("--output-root", type=Path, help="Output root for batch or library conversion.")
    parser.add_argument("--archive-source", action="store_true", help="Copy the original PPTX to source/original.pptx.")
    parser.add_argument("--tag-keywords", type=Path, help="Optional tag keyword file, either newline text or JSON list.")
    parser.add_argument("--no-resume", action="store_true", help="In library mode, reprocess decks even when manifest.json already exists.")
    parser.add_argument("--include-duplicates", action="store_true", help="In library mode, process duplicate SHA-256 files instead of cataloging them as duplicates.")
    parser.add_argument("--dry-run", action="store_true", help="In library mode, scan and write catalog/manifest without converting decks.")
    parser.add_argument("--exclude", action="append", default=[], help="In library mode, exclude paths by glob pattern. Can be repeated.")
    parser.add_argument("--standalone-html", action="store_true", help="Also write content.standalone.html with embedded image data for one-file sharing.")
    parser.add_argument("--rendered-html", action="store_true", help="Also render each slide as an image and write rendered.html for visual-fidelity browsing.")
    parser.add_argument("--rendered-standalone-html", action="store_true", help="Also render slides and write one self-contained rendered.standalone.html with embedded slide images.")
    parser.add_argument("--rendered-source-pdf", type=Path, help="Use an existing PDF as the visual render source instead of converting PPTX with soffice.")
    parser.add_argument("--progress-jsonl", action="store_true", help="Emit one JSON progress event per line while processing.")
    args = parser.parse_args()
    tag_keywords = load_tag_keywords(args.tag_keywords)
    if args.rendered_source_pdf and not args.pptx:
        parser.error("--rendered-source-pdf is only supported with --pptx")

    if args.pptx:
        if not args.output_dir:
            parser.error("--output-dir is required with --pptx")
        emit_progress(args.progress_jsonl, "progress_start", total=1, current=0, status="running", phase="convert")
        emit_progress(args.progress_jsonl, "file_start", total=1, current=1, status="running", file=str(args.pptx))
        result = process_deck(
            args.pptx,
            args.output_dir.resolve(),
            args.archive_source,
            tag_keywords,
            standalone_html=args.standalone_html,
            rendered_html=args.rendered_html,
            rendered_standalone_html=args.rendered_standalone_html,
            rendered_source_pdf=args.rendered_source_pdf,
        )
        emit_progress(args.progress_jsonl, "file_done", total=1, current=1, status="processed", file=str(args.pptx), output_dir=result["output_dir"])
        emit_progress(args.progress_jsonl, "progress_complete", total=1, current=1, status="done", phase="convert")
        print(json.dumps(result, ensure_ascii=False))
        return

    if args.library_root:
        if not args.output_root:
            parser.error("--output-root is required with --library-root")
        result = process_library(
            args.library_root,
            args.output_root,
            archive_source=args.archive_source,
            tag_keywords=tag_keywords,
            resume=not args.no_resume,
            skip_duplicates=not args.include_duplicates,
            dry_run=args.dry_run,
            exclude_patterns=args.exclude,
            standalone_html=args.standalone_html,
            rendered_html=args.rendered_html,
            rendered_standalone_html=args.rendered_standalone_html,
            progress_jsonl=args.progress_jsonl,
        )
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    if not args.output_root:
        parser.error("--output-root is required with --input-dir")

    input_dir = args.input_dir.resolve()
    if not input_dir.is_dir():
        raise NotADirectoryError(input_dir)
    results = []
    pptx_files = sorted(input_dir.glob("*.pptx"))
    emit_progress(args.progress_jsonl, "progress_start", total=len(pptx_files), current=0, status="running", phase="convert")
    for index, pptx_path in enumerate(pptx_files, 1):
        emit_progress(args.progress_jsonl, "file_start", total=len(pptx_files), current=index, status="running", file=str(pptx_path))
        output_dir = args.output_root.resolve() / safe_name(pptx_path.stem)
        try:
            results.append(
                process_deck(
                    pptx_path,
                    output_dir,
                    args.archive_source,
                    tag_keywords,
                    standalone_html=args.standalone_html,
                    rendered_html=args.rendered_html,
                    rendered_standalone_html=args.rendered_standalone_html,
                )
            )
            emit_progress(args.progress_jsonl, "file_done", total=len(pptx_files), current=index, status="processed", file=str(pptx_path), output_dir=str(output_dir))
        except Exception as exc:
            results.append({"input": str(pptx_path), "error": str(exc)})
            emit_progress(args.progress_jsonl, "file_done", total=len(pptx_files), current=index, status="failed", file=str(pptx_path), error=str(exc))
    emit_progress(args.progress_jsonl, "progress_complete", total=len(pptx_files), current=len(pptx_files), status="done", phase="convert")
    print(json.dumps({"deck_count": len(results), "results": results}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
